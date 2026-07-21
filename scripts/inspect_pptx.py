from __future__ import annotations

import argparse
import json
import sys
import zipfile
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")


def manifest_texts(manifest: dict[str, Any]) -> list[str]:
    values: list[str] = []
    for slide in manifest.get("slides", []):
        for item in slide.get("elements", []):
            text = item.get("text")
            if isinstance(text, str) and text.strip():
                values.append(text.strip())
    return values


def inspect(pptx_path: Path, manifest_path: Path | None = None) -> dict[str, Any]:
    prs = Presentation(pptx_path)
    slide_texts: list[str] = []
    picture_names: list[str] = []
    out_of_bounds: list[str] = []
    shape_count = 0
    connector_count = 0

    for slide in prs.slides:
        for shape in slide.shapes:
            shape_count += 1
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture_names.append(shape.name)
            if shape.shape_type == MSO_SHAPE_TYPE.LINE:
                connector_count += 1
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
            if (
                shape.left < 0
                or shape.top < 0
                or shape.left + shape.width > prs.slide_width + 10
                or shape.top + shape.height > prs.slide_height + 10
            ):
                out_of_bounds.append(shape.name)

    with zipfile.ZipFile(pptx_path) as archive:
        media_files = [
            name for name in archive.namelist()
            if name.startswith("ppt/media/") and not name.endswith("/")
        ]

    required_text: list[str] = []
    if manifest_path:
        manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
        required_text = manifest_texts(manifest)
    missing_text = [value for value in required_text if value not in slide_texts]

    result = {
        "pptx": str(pptx_path),
        "slides": len(prs.slides),
        "shapes": shape_count,
        "text_shapes_with_content": len(slide_texts),
        "connectors": connector_count,
        "picture_shapes": picture_names,
        "media_files": media_files,
        "missing_required_text": missing_text,
        "out_of_bounds_shapes": out_of_bounds,
        "passed": not picture_names and not media_files and not missing_text and not out_of_bounds,
    }
    return result


def main() -> int:
    parser = argparse.ArgumentParser(description="Inspect native editability and structure of a diagram PPTX.")
    parser.add_argument("--pptx", required=True, type=Path)
    parser.add_argument("--manifest", type=Path)
    parser.add_argument("--output", type=Path)
    args = parser.parse_args()

    result = inspect(args.pptx, args.manifest)
    if args.output:
        args.output.parent.mkdir(parents=True, exist_ok=True)
        args.output.write_text(json.dumps(result, ensure_ascii=False, indent=2), encoding="utf-8")
    print(json.dumps(result, ensure_ascii=False, indent=2))
    return 0 if result["passed"] else 2


if __name__ == "__main__":
    raise SystemExit(main())
