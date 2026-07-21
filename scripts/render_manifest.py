from __future__ import annotations

import argparse
import json
import sys
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


DEFAULT_FONT = "Microsoft YaHei"

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")


@dataclass(frozen=True)
class SlideScale:
    canvas_width: float
    canvas_height: float
    slide_width_in: float
    slide_height_in: float

    def x(self, value: float) -> int:
        return int(Inches(self.slide_width_in * float(value) / self.canvas_width))

    def y(self, value: float) -> int:
        return int(Inches(self.slide_height_in * float(value) / self.canvas_height))

    def w(self, value: float) -> int:
        return int(Inches(self.slide_width_in * float(value) / self.canvas_width))

    def h(self, value: float) -> int:
        return int(Inches(self.slide_height_in * float(value) / self.canvas_height))


SHAPE_TYPES = {
    "rect": MSO_AUTO_SHAPE_TYPE.RECTANGLE,
    "rectangle": MSO_AUTO_SHAPE_TYPE.RECTANGLE,
    "roundrect": MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
    "roundedrect": MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
    "ellipse": MSO_AUTO_SHAPE_TYPE.OVAL,
    "circle": MSO_AUTO_SHAPE_TYPE.OVAL,
    "diamond": MSO_AUTO_SHAPE_TYPE.DIAMOND,
    "hexagon": MSO_AUTO_SHAPE_TYPE.HEXAGON,
    "chevron": MSO_AUTO_SHAPE_TYPE.CHEVRON,
    "parallelogram": MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM,
    "cylinder": MSO_AUTO_SHAPE_TYPE.CAN,
    "cube": MSO_AUTO_SHAPE_TYPE.CUBE,
    "cloud": MSO_AUTO_SHAPE_TYPE.CLOUD,
    "terminator": MSO_AUTO_SHAPE_TYPE.FLOWCHART_TERMINATOR,
}

PROHIBITED_TYPES = {"image", "picture", "photo", "raster", "background_image"}


def parse_color(value: str | None, default: str = "#000000") -> RGBColor:
    text = (value or default).strip().lstrip("#")
    if len(text) == 3:
        text = "".join(ch * 2 for ch in text)
    if len(text) != 6:
        raise ValueError(f"Invalid color: {value!r}")
    return RGBColor(int(text[0:2], 16), int(text[2:4], 16), int(text[4:6], 16))


def suppress_theme_effects(shape: Any) -> None:
    element = shape._element
    style = element.find(qn("p:style"))
    if style is not None:
        effect_ref = style.find(qn("a:effectRef"))
        if effect_ref is not None:
            effect_ref.set("idx", "0")
    sp_pr = element.find(qn("p:spPr"))
    if sp_pr is not None:
        for node in list(sp_pr):
            if node.tag in {qn("a:effectLst"), qn("a:effectDag")}:
                sp_pr.remove(node)
        sp_pr.append(OxmlElement("a:effectLst"))


def set_fill_transparency(shape: Any, value: float) -> None:
    if value <= 0:
        return
    solid = shape.fill._xPr.find(qn("a:solidFill"))
    if solid is None:
        return
    color_node = next(iter(solid), None)
    if color_node is None:
        return
    alpha = OxmlElement("a:alpha")
    alpha.set("val", str(max(0, min(100000, int((100 - value) * 1000)))))
    color_node.append(alpha)


def set_east_asia_font(run: Any, font_name: str) -> None:
    run.font.name = font_name
    rpr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        node = rpr.find(qn(tag))
        if node is None:
            node = OxmlElement(tag)
            rpr.append(node)
        node.set("typeface", font_name)


def set_dash(line: Any, dash: bool) -> None:
    if not dash:
        return
    ln = line._get_or_add_ln()
    for node in list(ln):
        if node.tag == qn("a:prstDash"):
            ln.remove(node)
    node = OxmlElement("a:prstDash")
    node.set("val", "dash")
    ln.append(node)


def set_arrowheads(line: Any, arrow_start: bool, arrow_end: bool) -> None:
    ln = line._get_or_add_ln()
    for tag, enabled in (("a:headEnd", arrow_start), ("a:tailEnd", arrow_end)):
        for node in list(ln):
            if node.tag == qn(tag):
                ln.remove(node)
        node = OxmlElement(tag)
        node.set("type", "triangle" if enabled else "none")
        node.set("w", "sm")
        node.set("len", "sm")
        ln.append(node)


def alignment(value: str | None) -> PP_ALIGN:
    return {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }.get((value or "left").lower(), PP_ALIGN.LEFT)


def vertical_alignment(value: str | None) -> MSO_ANCHOR:
    return {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "center": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }.get((value or "middle").lower(), MSO_ANCHOR.MIDDLE)


def format_text_frame(frame: Any, item: dict[str, Any]) -> None:
    frame.clear()
    frame.word_wrap = bool(item.get("wrap", True))
    margin = float(item.get("margin_px", 0))
    frame.margin_left = frame.margin_right = int(Inches(margin / 96.0))
    frame.margin_top = frame.margin_bottom = int(Inches(margin / 96.0))
    frame.vertical_anchor = vertical_alignment(item.get("valign"))
    paragraph = frame.paragraphs[0]
    paragraph.alignment = alignment(item.get("align"))
    paragraph.line_spacing = float(item.get("line_spacing", 1.0))
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)
    run = paragraph.add_run()
    run.text = str(item.get("text", ""))
    set_east_asia_font(run, str(item.get("font_family") or DEFAULT_FONT))
    run.font.size = Pt(float(item.get("font_size_pt", 14)))
    run.font.bold = bool(item.get("bold", False))
    run.font.italic = bool(item.get("italic", False))
    run.font.color.rgb = parse_color(item.get("color"), "#111827")


def add_text(slide: Any, item: dict[str, Any], scale: SlideScale) -> Any:
    shape = slide.shapes.add_textbox(
        scale.x(item["x"]), scale.y(item["y"]), scale.w(item["w"]), scale.h(item["h"])
    )
    shape.name = str(item.get("name") or item.get("id") or "Text")
    format_text_frame(shape.text_frame, item)
    return shape


def add_shape(slide: Any, item: dict[str, Any], scale: SlideScale) -> Any:
    key = str(item.get("shape", "rect")).replace("_", "").lower()
    if key not in SHAPE_TYPES:
        raise ValueError(f"Unsupported shape: {item.get('shape')!r}")
    shape = slide.shapes.add_shape(
        SHAPE_TYPES[key],
        scale.x(item["x"]), scale.y(item["y"]), scale.w(item["w"]), scale.h(item["h"]),
    )
    shape.name = str(item.get("name") or item.get("id") or key)
    fill = str(item.get("fill", "#FFFFFF"))
    if fill.lower() in {"none", "transparent"}:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = parse_color(fill, "#FFFFFF")
        set_fill_transparency(shape, float(item.get("fill_transparency", 0)))
    stroke = str(item.get("stroke", "#64748B"))
    if stroke.lower() in {"none", "transparent"}:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = parse_color(stroke, "#64748B")
        shape.line.width = Pt(float(item.get("stroke_width_pt", 1.0)))
        set_dash(shape.line, bool(item.get("dash", False)))
    suppress_theme_effects(shape)
    if "text" in item:
        format_text_frame(shape.text_frame, item)
    return shape


def add_segment(
    slide: Any,
    scale: SlideScale,
    start: list[float],
    end: list[float],
    item: dict[str, Any],
    arrow_start: bool,
    arrow_end: bool,
    name: str,
) -> Any:
    shape = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        scale.x(start[0]), scale.y(start[1]), scale.x(end[0]), scale.y(end[1]),
    )
    shape.name = name
    shape.line.color.rgb = parse_color(item.get("stroke"), "#334155")
    shape.line.width = Pt(float(item.get("stroke_width_pt", 1.5)))
    set_dash(shape.line, bool(item.get("dash", False)))
    set_arrowheads(shape.line, arrow_start, arrow_end)
    suppress_theme_effects(shape)
    return shape


def add_connector(slide: Any, item: dict[str, Any], scale: SlideScale) -> int:
    points = item.get("points")
    if not isinstance(points, list) or len(points) < 2:
        raise ValueError(f"Connector {item.get('id')!r} requires at least two points")
    base_name = str(item.get("name") or item.get("id") or "Connector")
    for index, (start, end) in enumerate(zip(points, points[1:])):
        add_segment(
            slide,
            scale,
            start,
            end,
            item,
            bool(item.get("arrow_start", False)) and index == 0,
            bool(item.get("arrow_end", True)) and index == len(points) - 2,
            f"{base_name}_{index + 1}",
        )
    return len(points) - 1


def symbol_shape(
    slide: Any,
    scale: SlideScale,
    shape_type: MSO_AUTO_SHAPE_TYPE,
    x: float,
    y: float,
    w: float,
    h: float,
    color: str,
    name: str,
    fill: str = "none",
    stroke_width_pt: float = 1.4,
) -> Any:
    shape = slide.shapes.add_shape(shape_type, scale.x(x), scale.y(y), scale.w(w), scale.h(h))
    shape.name = name
    if fill.lower() in {"none", "transparent"}:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = parse_color(fill)
    shape.line.color.rgb = parse_color(color)
    shape.line.width = Pt(stroke_width_pt)
    suppress_theme_effects(shape)
    return shape


def symbol_line(
    slide: Any,
    scale: SlideScale,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    color: str,
    name: str,
    width_pt: float = 1.2,
    arrow_end: bool = False,
) -> Any:
    return add_segment(
        slide, scale, [x1, y1], [x2, y2],
        {"stroke": color, "stroke_width_pt": width_pt, "dash": False}, False, arrow_end, name,
    )


def add_symbol(slide: Any, item: dict[str, Any], scale: SlideScale) -> int:
    kind = str(item.get("symbol", "")).lower()
    x, y, w, h = map(float, (item["x"], item["y"], item["w"], item["h"]))
    color = str(item.get("color", "#1677FF"))
    name = str(item.get("name") or item.get("id") or kind)
    count = 0
    base = 60.0
    factor = min(w / base, h / base)
    ox = x + (w - base * factor) / 2
    oy = y + (h - base * factor) / 2

    def px(value: float) -> float:
        return ox + value * factor

    def py(value: float) -> float:
        return oy + value * factor

    def pw(value: float) -> float:
        return value * factor

    def ph(value: float) -> float:
        return value * factor

    def shape(
        shape_type: MSO_AUTO_SHAPE_TYPE,
        rx: float,
        ry: float,
        rw: float,
        rh: float,
        suffix: str,
        fill: str = "none",
        width: float = 1.4,
    ) -> None:
        nonlocal count
        symbol_shape(
            slide, scale, shape_type, px(rx), py(ry), pw(rw), ph(rh), color,
            f"{name}_{suffix}", fill=fill, stroke_width_pt=width,
        )
        count += 1

    def line(
        x1: float,
        y1: float,
        x2: float,
        y2: float,
        suffix: str,
        width: float = 1.2,
        arrow_end: bool = False,
        line_color: str | None = None,
    ) -> None:
        nonlocal count
        symbol_line(
            slide, scale, px(x1), py(y1), px(x2), py(y2), line_color or color,
            f"{name}_{suffix}", width_pt=width, arrow_end=arrow_end,
        )
        count += 1

    def polyline(points: list[tuple[float, float]], suffix: str, width: float = 1.2, arrow_end: bool = False, line_color: str | None = None) -> None:
        for index, (start, end) in enumerate(zip(points, points[1:])):
            line(
                start[0], start[1], end[0], end[1], f"{suffix}_{index + 1}", width,
                arrow_end=arrow_end and index == len(points) - 2, line_color=line_color,
            )

    if kind == "database":
        shape(MSO_AUTO_SHAPE_TYPE.CAN, 8, 5, 44, 50, "body", width=1.5)
        line(14, 25, 46, 25, "layer_1", 1.0)
        line(14, 39, 46, 39, "layer_2", 1.0)
        return count
    if kind == "cube":
        shape(MSO_AUTO_SHAPE_TYPE.CUBE, 5, 5, 50, 50, "body", width=1.6)
        return count
    if kind == "gear":
        shape(MSO_AUTO_SHAPE_TYPE.GEAR_6, 5, 5, 50, 50, "body", width=1.4)
        return count
    if kind == "person":
        shape(MSO_AUTO_SHAPE_TYPE.OVAL, 23, 2, 14, 14, "head", width=1.5)
        line(30, 16, 30, 43, "body", 1.6)
        line(30, 24, 13, 34, "arm_1", 1.6)
        line(30, 24, 47, 34, "arm_2", 1.6)
        line(30, 43, 17, 58, "leg_1", 1.6)
        line(30, 43, 43, 58, "leg_2", 1.6)
        return count
    if kind == "people":
        for index, (hx, hy, diameter) in enumerate(((25, 5, 10), (8, 18, 8), (44, 18, 8))):
            shape(MSO_AUTO_SHAPE_TYPE.OVAL, hx, hy, diameter, diameter, f"head_{index + 1}", fill="#FFFFFF", width=1.5)
        line(30, 16, 30, 39, "body", 1.6)
        line(30, 23, 16, 33, "arm_1", 1.6)
        line(30, 23, 44, 33, "arm_2", 1.6)
        line(30, 39, 19, 53, "leg_1", 1.6)
        line(30, 39, 41, 53, "leg_2", 1.6)
        return count
    if kind == "monitor":
        shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 5, 5, 50, 37, "screen", width=1.6)
        polyline([(11,26),(18,26),(22,15),(27,32),(33,21),(40,21),(47,17)], "wave", 1.3)
        line(25, 42, 25, 51, "stand_1", 1.5)
        line(35, 42, 35, 51, "stand_2", 1.5)
        line(18, 51, 42, 51, "base", 1.5)
        return count
    if kind in {"clipboard", "clipboard-check"}:
        shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 8, 8, 44, 49, "body", width=1.5)
        shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 20, 3, 20, 10, "clip", fill="#FFFFFF", width=1.4)
        for index, cy in enumerate((24, 35, 46)):
            shape(MSO_AUTO_SHAPE_TYPE.OVAL, 15, cy - 2, 4, 4, f"bullet_{index + 1}", fill=color, width=0.5)
            line(24, cy, 43, cy, f"text_{index + 1}", 1.0)
        if kind == "clipboard-check":
            shape(MSO_AUTO_SHAPE_TYPE.OVAL, 37, 40, 18, 18, "status", fill="#FFFFFF", width=1.2)
            polyline([(41,49),(45,53),(52,44)], "check", 1.2)
        return count
    if kind in {"document", "server"}:
        shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 8, 5, 44, 50, "body", width=1.4)
        for index, yy in enumerate((20, 31, 42)):
            line(16, yy, 44, yy, f"line_{index + 1}", 1.0)
        return count
    if kind == "report":
        shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 8, 5, 44, 50, "document", width=1.5)
        for index, yy in enumerate((16, 25, 34)):
            shape(MSO_AUTO_SHAPE_TYPE.OVAL, 14, yy, 3, 3, f"dot_{index + 1}", fill=color, width=0.5)
            line(21, yy + 1.5, 38, yy + 1.5, f"text_{index + 1}", 0.9)
        for index, (bx, height) in enumerate(((31,8),(38,14),(45,20))):
            shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, bx, 50-height, 4, height, f"bar_{index + 1}", fill=color, width=0.5)
        return count
    if kind == "warehouse":
        polyline([(3,22),(30,4),(57,22)], "roof", 1.8)
        shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 8, 22, 44, 34, "body", width=1.5)
        shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 22, 35, 17, 21, "door", width=1.2)
        shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 12, 28, 9, 8, "box_1", width=1.0)
        shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 40, 28, 9, 8, "box_2", width=1.0)
        return count
    if kind == "book":
        shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 5, 8, 25, 44, "left_page", width=1.4)
        shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 30, 8, 25, 44, "right_page", width=1.4)
        line(30, 12, 30, 49, "spine", 1.1)
        for index, yy in enumerate((21, 31, 41)):
            line(10, yy, 23, yy, f"left_line_{index + 1}", 0.9)
            line(37, yy, 50, yy, f"right_line_{index + 1}", 0.9)
        return count
    if kind == "arm":
        shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 5, 48, 42, 8, "base", width=1.4)
        joints = [(13,42),(27,25),(42,12),(53,25)]
        for index, (jx, jy) in enumerate(joints):
            shape(MSO_AUTO_SHAPE_TYPE.OVAL, jx-4, jy-4, 8, 8, f"joint_{index + 1}", fill="#FFFFFF", width=1.5)
        for index, (start, end) in enumerate(zip(joints, joints[1:])):
            line(start[0], start[1], end[0], end[1], f"arm_{index + 1}", 2.0)
        line(53, 25, 59, 32, "claw_1", 1.7)
        line(53, 25, 59, 19, "claw_2", 1.7)
        return count
    if kind == "brain":
        for index, (cx, cy) in enumerate(((8,10),(20,6),(33,11),(8,27),(21,27),(35,28))):
            shape(MSO_AUTO_SHAPE_TYPE.OVAL, cx, cy, 16, 16, f"lobe_{index + 1}", fill=color, width=0.8)
        line(29, 9, 29, 47, "center", 1.2, line_color="#FFFFFF")
        for index, (x1,y1,x2,y2) in enumerate(((15,18,25,18),(34,17,43,21),(15,36,25,31),(34,35,44,30))):
            line(x1, y1, x2, y2, f"fold_{index + 1}", 1.0, line_color="#FFFFFF")
        return count
    if kind == "tools":
        shape(MSO_AUTO_SHAPE_TYPE.GEAR_6, 27, 24, 29, 29, "gear", width=1.3)
        line(12, 10, 38, 36, "handle", 2.0)
        shape(MSO_AUTO_SHAPE_TYPE.OVAL, 5, 3, 14, 14, "head", fill="#FFFFFF", width=1.5)
        return count
    if kind == "bars":
        for index, height in enumerate((15,29,43)):
            shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 7 + index*14, 51-height, 8, height, f"bar_{index + 1}", width=1.2)
        line(3, 51, 49, 51, "axis", 1.2)
        return count
    if kind == "refresh":
        shape(MSO_AUTO_SHAPE_TYPE.OVAL, 9, 9, 42, 42, "ring", width=1.5)
        polyline([(39,12),(51,12),(51,24)], "arrow_1", 1.5, arrow_end=True)
        polyline([(21,48),(9,48),(9,36)], "arrow_2", 1.5, arrow_end=True)
        return count
    raise ValueError(f"Unsupported symbol: {kind!r}")


def render(manifest: dict[str, Any], output: Path) -> dict[str, Any]:
    deck = manifest.get("deck") or {}
    canvas_width = float(deck.get("canvas_width", 1600))
    canvas_height = float(deck.get("canvas_height", 900))
    slide_width_in = float(deck.get("slide_width_in", 13.333))
    slide_height_in = float(deck.get("slide_height_in", slide_width_in * canvas_height / canvas_width))
    scale = SlideScale(canvas_width, canvas_height, slide_width_in, slide_height_in)

    slides_data = manifest.get("slides")
    if not isinstance(slides_data, list) or not slides_data:
        raise ValueError("Manifest requires a non-empty slides list")

    prs = Presentation()
    prs.slide_width = Inches(slide_width_in)
    prs.slide_height = Inches(slide_height_in)
    counters = {"slides": 0, "elements": 0, "text": 0, "shape": 0, "connector_segments": 0, "symbol_shapes": 0}

    for slide_data in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.name = str(slide_data.get("name") or f"Slide {len(prs.slides)}")
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = parse_color(slide_data.get("background"), "#FFFFFF")
        counters["slides"] += 1
        for item in slide_data.get("elements", []):
            element_type = str(item.get("type", "")).lower()
            if element_type in PROHIBITED_TYPES:
                raise ValueError(f"Raster element type is prohibited: {element_type}")
            if element_type == "text":
                add_text(slide, item, scale)
                counters["text"] += 1
            elif element_type == "shape":
                add_shape(slide, item, scale)
                counters["shape"] += 1
                if "text" in item:
                    counters["text"] += 1
            elif element_type == "connector":
                counters["connector_segments"] += add_connector(slide, item, scale)
            elif element_type == "symbol":
                counters["symbol_shapes"] += add_symbol(slide, item, scale)
            else:
                raise ValueError(f"Unsupported element type: {element_type!r}")
            counters["elements"] += 1

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    counters["output"] = str(output)
    counters["canvas"] = [canvas_width, canvas_height]
    counters["slide_size_inches"] = [slide_width_in, slide_height_in]
    return counters


def main() -> int:
    parser = argparse.ArgumentParser(description="Render an editable PowerPoint diagram from a JSON manifest.")
    parser.add_argument("--manifest", required=True, type=Path)
    parser.add_argument("--output", required=True, type=Path)
    parser.add_argument("--summary", type=Path)
    args = parser.parse_args()

    manifest = json.loads(args.manifest.read_text(encoding="utf-8"))
    summary = render(manifest, args.output)
    if args.summary:
        args.summary.parent.mkdir(parents=True, exist_ok=True)
        args.summary.write_text(json.dumps(summary, ensure_ascii=False, indent=2), encoding="utf-8")
    print(json.dumps(summary, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
